#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
import os
import tempfile
import uuid
from io import BytesIO

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Optional: Pillow for robust image re-encode (skips if not available)
try:
    from PIL import Image
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


# -----------------------------
# Helpers
# -----------------------------
VALID_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp"}

def ensure_dir(path: str):
    if not path:
        return
    os.makedirs(path, exist_ok=True)

def read_csv_safe(p: str) -> pd.DataFrame:
    return pd.read_csv(p, dtype=str, low_memory=False)

def to_float(s: pd.Series) -> pd.Series:
    return pd.to_numeric(s, errors="coerce")

def pct_fmt(x):
    if pd.isna(x):
        return "-"
    return f"{x:.0f}%"

def try_url(u: str) -> str | None:
    if not isinstance(u, str):
        return None
    u = u.strip()
    if u.lower().startswith(("http://", "https://")):
        return u
    return None

def clean_path(p: str) -> str | None:
    if not isinstance(p, str):
        return None
    s = p.strip().strip('"').strip("'").replace("/", "\\")
    return s

def is_local_img(p: str) -> bool:
    if not isinstance(p, str):
        return False
    ext = os.path.splitext(p)[1].lower()
    return os.path.isfile(p) and ext in VALID_IMAGE_EXTS

def sanitize_image_to_png(img_path: str, tmp_dir: str) -> str | None:
    """
    Return a safe PNG path for pptx, or None if cannot be sanitized.
    If Pillow is unavailable, return img_path if it already exists and is a valid image.
    """
    if not is_local_img(img_path):
        return None
    if not PIL_AVAILABLE:
        # No Pillow: return original (pptx can handle common formats)
        return img_path
    try:
        os.makedirs(tmp_dir, exist_ok=True)
        out_png = os.path.join(tmp_dir, f"{uuid.uuid4().hex}.png")
        im = Image.open(img_path).convert("RGB")
        im.save(out_png, format="PNG")
        return out_png
    except Exception:
        return None

def safe_add_picture(slide, img_path, left_in, top_in, width_in=None, height_in=None):
    """
    Add a picture safely. Re-encodes to PNG in a temp dir if Pillow is available.
    Returns shape or None if skipped.
    """
    tmp_dir = os.path.join(tempfile.gettempdir(), "pptx_imgs")
    safe = sanitize_image_to_png(img_path, tmp_dir)
    if not safe:
        return None
    try:
        if width_in is not None:
            return slide.shapes.add_picture(safe, Inches(left_in), Inches(top_in), width=Inches(width_in))
        elif height_in is not None:
            return slide.shapes.add_picture(safe, Inches(left_in), Inches(top_in), height=Inches(height_in))
        else:
            return slide.shapes.add_picture(safe, Inches(left_in), Inches(top_in))
    except Exception:
        return None

def fig_to_picture(slide, fig, left_in, top_in, width_in):
    """
    Save matplotlib fig to a real PNG file (not BytesIO) to avoid Windows packaging issues,
    then place it on the slide.
    """
    tmp_dir = os.path.join(tempfile.gettempdir(), "pptx_figs")
    os.makedirs(tmp_dir, exist_ok=True)
    tmp_path = os.path.join(tmp_dir, f"{uuid.uuid4().hex}.png")
    fig.savefig(tmp_path, format="png", dpi=180, bbox_inches="tight")
    plt.close(fig)
    slide.shapes.add_picture(tmp_path, Inches(left_in), Inches(top_in), width=Inches(width_in))

def add_title(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title box
    left, top, width, height = Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.6))
        tf2 = tx2.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.LEFT
    return slide

def add_kv_block(slide, rows, left=0.5, top=1.6, col1w=2.5, col2w=6.0, row_h=0.35):
    """
    Simple key-value text blocks stacked vertically.
    rows: list of (label, value) strings
    """
    y = top
    for (k, v) in rows:
        # key
        tb1 = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(col1w), Inches(row_h))
        tf1 = tb1.text_frame; tf1.clear()
        p1 = tf1.paragraphs[0]; r1 = p1.add_run()
        r1.text = str(k); r1.font.size = Pt(14); r1.font.bold = True

        # value
        tb2 = slide.shapes.add_textbox(Inches(left + col1w + 0.2), Inches(y), Inches(col2w), Inches(row_h))
        tf2 = tb2.text_frame; tf2.clear()
        p2 = tf2.paragraphs[0]; r2 = p2.add_run()
        r2.text = str(v); r2.font.size = Pt(14)
        y += row_h + 0.06

def add_table(slide, df: pd.DataFrame, left=0.5, top=1.3, col_w=1.4, row_h=0.28, max_rows=20):
    """
    Add a compact table from a DataFrame.
    """
    data = df.copy()
    if len(data) > max_rows:
        data = data.head(max_rows)
    rows, cols = data.shape
    table = slide.shapes.add_table(rows+1, cols, Inches(left), Inches(top), Inches(col_w*cols), Inches(row_h*(rows+1))).table

    # header
    for j, c in enumerate(data.columns):
        cell = table.cell(0, j)
        cell.text = str(c)
        cell.text_frame.paragraphs[0].runs[0].font.bold = True

    # rows
    for i in range(rows):
        for j in range(cols):
            val = data.iat[i, j]
            cell = table.cell(i+1, j)
            cell.text = "" if pd.isna(val) else str(val)

def top_bottom(df, n=5, score_col="composite_score"):
    s = to_float(df.get(score_col, pd.Series(dtype="float64")))
    df2 = df.copy()
    df2["_score_float"] = s
    df2 = df2.dropna(subset=["_score_float"])
    df2 = df2.sort_values("_score_float", ascending=False)
    topn = df2.head(n)
    botn = df2.tail(n).sort_values("_score_float", ascending=True)
    return topn, botn

def load_covers(covers_path: str | None) -> dict:
    """
    Returns a dict: {video_id: local_image_path}
    Supports flexible column names: 'cover_path', 'frame_file', 'path', 'thumb', 'image'
    """
    if not covers_path or not os.path.isfile(covers_path):
        return {}
    cov = read_csv_safe(covers_path)
    # Normalize column names a bit
    cols = [c.lower() for c in cov.columns]
    cov.columns = cols
    vid_col = "video_id" if "video_id" in cols else None
    if not vid_col:
        return {}

    path_col = None
    for cname in ["cover_path", "frame_file", "path", "thumb", "image", "file"]:
        if cname in cols:
            path_col = cname
            break
    if not path_col:
        return {}

    cov["__p__"] = cov[path_col].map(clean_path)
    cov = cov.dropna(subset=["__p__", "video_id"])
    # Keep only existing local images
    cov = cov[cov["__p__"].map(is_local_img)]
    mapping = dict(zip(cov["video_id"], cov["__p__"]))
    return mapping


# -----------------------------
# Slide Builders
# -----------------------------
def overview_slide(prs, sc: pd.DataFrame):
    slide = add_title(prs, "TikTok Executive Summary")

    # Basic metrics
    n_posts = len(sc)
    score = to_float(sc.get("composite_score", pd.Series(dtype="float64")))
    avg_score = float(np.nanmean(score)) if len(score) else float("nan")

    # brand tone and overlay coverage
    brand_pct = to_float(sc.get("brand_match_frame_percent", pd.Series(dtype="float64")))
    overlay_pct = to_float(sc.get("overlay_frame_percent", pd.Series(dtype="float64")))
    cuts_per_min = to_float(sc.get("cuts_per_minute", pd.Series(dtype="float64")))

    # disclosure rate
    disc = sc.get("disclosure_ok")
    if disc is not None:
        disc_rate = (disc.astype(str).str.lower().isin(["1", "true", "yes"])).mean() * 100
    else:
        disc_rate = float("nan")

    rows = [
        ("Posts analyzed", n_posts),
        ("Average score", f"{avg_score:.1f}" if pd.notna(avg_score) else "-"),
        ("Avg brand tone coverage", pct_fmt(np.nanmean(brand_pct)) if len(brand_pct) else "-"),
        ("Avg overlay coverage", pct_fmt(np.nanmean(overlay_pct)) if len(overlay_pct) else "-"),
        ("Avg cuts/min", f"{np.nanmean(cuts_per_min):.2f}" if len(cuts_per_min) else "-"),
        ("Disclosure rate", pct_fmt(disc_rate)),
    ]
    add_kv_block(slide, rows)

    # Charts: Score histogram & brand vs overlay scatter (if data present)
    # Histogram
    if len(score.dropna()):
        fig = plt.figure(figsize=(8, 3.2))
        plt.hist(score.dropna().values, bins=20)
        plt.title("Score Distribution")
        plt.xlabel("Composite Score")
        plt.ylabel("Count")
        fig_to_picture(slide, fig, left_in=0.6, top_in=3.2, width_in=8.8)

    # Scatter brand% vs overlay%
    if len(brand_pct.dropna()) and len(overlay_pct.dropna()):
        df_sc = pd.DataFrame({
            "brand_pct": brand_pct,
            "overlay_pct": overlay_pct
        }).dropna()
        if len(df_sc) > 0:
            fig = plt.figure(figsize=(8, 3.2))
            plt.scatter(df_sc["brand_pct"], df_sc["overlay_pct"], s=10)
            plt.title("Brand Tone % vs Overlay %")
            plt.xlabel("Brand tone coverage %")
            plt.ylabel("Overlay coverage %")
            fig_to_picture(slide, fig, left_in=0.6, top_in=6.0, width_in=8.8)

def by_brand_slide(prs, sc: pd.DataFrame):
    slide = add_title(prs, "By Brand")

    # derive brand column if missing
    brand = sc.get("brand")
    if brand is None:
        brand = sc.get("handle", pd.Series(dtype="object"))
    sc2 = sc.copy()
    sc2["__brand__"] = brand

    # compute pivots
    score = to_float(sc2.get("composite_score", pd.Series(dtype="float64")))
    brand_pct = to_float(sc2.get("brand_match_frame_percent", pd.Series(dtype="float64")))
    overlay_pct = to_float(sc2.get("overlay_frame_percent", pd.Series(dtype="float64")))
    cuts = to_float(sc2.get("cuts_per_minute", pd.Series(dtype="float64")))
    disc = sc2.get("disclosure_ok")
    disc_rate = (disc.astype(str).str.lower().isin(["1","true","yes"])) if disc is not None else pd.Series([np.nan]*len(sc2))

    df = pd.DataFrame({
        "brand": sc2["__brand__"],
        "score": score,
        "brand_pct": brand_pct,
        "overlay_pct": overlay_pct,
        "cuts": cuts,
        "disc": disc_rate
    })

    # group safely
    grp = df.groupby("brand", dropna=False)
    tbl = pd.DataFrame({
        "posts": grp.size()
    })
    tbl["avg_score"] = grp["score"].mean()
    tbl["brand_pct"] = grp["brand_pct"].mean()
    tbl["overlay_pct"] = grp["overlay_pct"].mean()
    tbl["cuts_min"] = grp["cuts"].mean()
    tbl["disclosure_rate"] = grp["disc"].mean() * 100

    # tidy
    tbl = tbl.reset_index().rename(columns={"brand": "Brand"})
    # format a display copy
    disp = tbl.copy()
    disp["avg_score"] = disp["avg_score"].map(lambda v: f"{v:.1f}" if pd.notna(v) else "-")
    disp["brand_pct"] = disp["brand_pct"].map(lambda v: pct_fmt(v))
    disp["overlay_pct"] = disp["overlay_pct"].map(lambda v: pct_fmt(v))
    disp["cuts_min"] = disp["cuts_min"].map(lambda v: f"{v:.2f}" if pd.notna(v) else "-")
    disp["disclosure_rate"] = disp["disclosure_rate"].map(lambda v: pct_fmt(v))

    add_table(slide, disp[["Brand","posts","avg_score","brand_pct","overlay_pct","cuts_min","disclosure_rate"]],
              left=0.5, top=1.4, col_w=1.5, row_h=0.3, max_rows=22)

def by_hotel_slide(prs, sc: pd.DataFrame):
    if "hotel" not in sc.columns:
        return  # silently skip
    slide = add_title(prs, "By Hotel")
    score = to_float(sc.get("composite_score", pd.Series(dtype="float64")))
    brand_pct = to_float(sc.get("brand_match_frame_percent", pd.Series(dtype="float64")))
    overlay_pct = to_float(sc.get("overlay_frame_percent", pd.Series(dtype="float64")))
    cuts = to_float(sc.get("cuts_per_minute", pd.Series(dtype="float64")))
    disc = sc.get("disclosure_ok")
    disc_rate = (disc.astype(str).str.lower().isin(["1","true","yes"])) if disc is not None else pd.Series([np.nan]*len(sc))

    df = pd.DataFrame({
        "hotel": sc["hotel"],
        "score": score,
        "brand_pct": brand_pct,
        "overlay_pct": overlay_pct,
        "cuts": cuts,
        "disc": disc_rate
    })
    grp = df.groupby("hotel", dropna=False)
    tbl = pd.DataFrame({
        "posts": grp.size()
    })
    tbl["avg_score"] = grp["score"].mean()
    tbl["brand_pct"] = grp["brand_pct"].mean()
    tbl["overlay_pct"] = grp["overlay_pct"].mean()
    tbl["cuts_min"] = grp["cuts"].mean()
    tbl["disclosure_rate"] = grp["disc"].mean() * 100
    tbl = tbl.reset_index().rename(columns={"hotel":"Hotel"})

    disp = tbl.copy()
    disp["avg_score"] = disp["avg_score"].map(lambda v: f"{v:.1f}" if pd.notna(v) else "-")
    disp["brand_pct"] = disp["brand_pct"].map(pct_fmt)
    disp["overlay_pct"] = disp["overlay_pct"].map(pct_fmt)
    disp["cuts_min"] = disp["cuts_min"].map(lambda v: f"{v:.2f}" if pd.notna(v) else "-")
    disp["disclosure_rate"] = disp["disclosure_rate"].map(pct_fmt)

    add_table(slide, disp[["Hotel","posts","avg_score","brand_pct","overlay_pct","cuts_min","disclosure_rate"]],
              left=0.5, top=1.4, col_w=1.5, row_h=0.3, max_rows=20)

def thumb_grid_slide(prs, title, rows_df: pd.DataFrame, covers_map: dict, n=5):
    slide = add_title(prs, title)
    # Columns we try to display
    cols_try = ["video_id","handle","brand","composite_score","video_url"]
    dd = {}
    for c in cols_try:
        if c in rows_df.columns:
            dd[c] = rows_df[c].values.tolist()
        else:
            dd[c] = [None]*len(rows_df)
    # layout: 5 rows, each row = thumbnail + text
    y = 1.4
    for i in range(min(n, len(rows_df))):
        vid = dd["video_id"][i]
        handle = dd["handle"][i]
        brand = dd["brand"][i]
        score = dd["composite_score"][i]
        url = try_url(dd["video_url"][i])

        # image on left, 1.5" high (width maintains proportion if available)
        img = covers_map.get(str(vid)) if vid is not None else None
        _ = safe_add_picture(slide, img, left_in=0.5, top_in=y, height_in=1.5)

        # text on right
        tb = slide.shapes.add_textbox(Inches(2.2), Inches(y), Inches(7.0), Inches(1.5))
        tf = tb.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = f"{i+1}. {handle or '-'}  |  {brand or '-'}"
        r0.font.size = Pt(14); r0.font.bold = True

        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = f"Score: {score}" if score is not None else "Score: -"
        r1.font.size = Pt(12)

        if url:
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = url
            r2.font.size = Pt(11)
        y += 1.6

def top_bottom_slides(prs, sc: pd.DataFrame, covers_map: dict):
    topn, botn = top_bottom(sc, n=5, score_col="composite_score")
    if len(topn):
        thumb_grid_slide(prs, "Top 5", topn, covers_map, n=5)
    if len(botn):
        thumb_grid_slide(prs, "Bottom 5", botn, covers_map, n=5)


# -----------------------------
# Main
# -----------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--scoring", required=True, help="Path to scoring_summary.csv")
    parser.add_argument("--covers", default=None, help="Optional covers manifest CSV (video_id -> cover_path)")
    parser.add_argument("--out", default="Exec_Summary.pptx", help="Output PPTX path")
    args = parser.parse_args()

    # Output dir exists?
    out_dir = os.path.dirname(args.out)
    if out_dir:
        ensure_dir(out_dir)

    # Load data
    sc = read_csv_safe(args.scoring)
    # Normalize column names (but keep original for display)
    # (We won't downcase globally to avoid breaking user column names elsewhere)
    covers_map = load_covers(args.covers)

    # Build deck
    prs = Presentation()

    # Slides
    overview_slide(prs, sc)
    by_brand_slide(prs, sc)
    by_hotel_slide(prs, sc)   # silently skips if no 'hotel' col
    top_bottom_slides(prs, sc, covers_map)

    # Save
    prs.save(args.out)
    print(f"Wrote: {args.out}")

if __name__ == "__main__":
    main()
