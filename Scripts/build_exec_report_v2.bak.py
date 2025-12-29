#!/usr/bin/env python3
# build_exec_report_v2.py
# Adds thumbnail previews (from covers_manifest) and brand-by-brand scorecards.

import argparse, pandas as pd, numpy as np, os, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt

# Optional: requests for downloading cover_url if no local saved_path exists
try:
    import requests
    HAVE_REQUESTS = True
except Exception:
    HAVE_REQUESTS = False

def read_csv_safe(p):
    return pd.read_csv(p, dtype=str, low_memory=False)

def to_float(s):
    try: return float(s)
    except: return np.nan

def ensure_brand_col(df):
    if "brand" in df.columns:
        return df
    # fallback brand from handle
    def brand_from_handle(h):
        if not isinstance(h, str): return "Other"
        x=h.lower()
        if "anantara" in x: return "Anantara"
        if "avani" in x: return "Avani"
        if "nh" in x and "collection" in x: return "NH Collection"
        if x.startswith("nh") or " nh" in x or "nhhotel" in x: return "NH"
        if "tivoli" in x: return "Tivoli"
        if "oaks" in x: return "Oaks"
        return "Other"
    df["brand"] = df.get("handle","").map(brand_from_handle)
    return df

def add_title(slide, text):
    slide.shapes.title.text = text

def add_kpi_box(slide, left_in, top_in, title, value):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(2.6), Inches(1.05))
    tf = box.text_frame; tf.clear()
    p1 = tf.paragraphs[0]; p1.text = title; p1.font.size = Pt(12)
    p2 = tf.add_paragraph(); p2.text = value; p2.font.size = Pt(24); p2.font.bold = True

def fig_to_picture(slide, fig, left_in, top_in, width_in):
    bio = io.BytesIO()
    fig.savefig(bio, format="png", dpi=180, bbox_inches="tight")
    bio.seek(0)
    slide.shapes.add_picture(bio, Inches(left_in), Inches(top_in), width=Inches(width_in))
    plt.close(fig)

def brand_mix_chart(sc, prs):
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(s1, "TikTok Creative Effectiveness — Executive Summary")

    # KPI tiles
    sc["brand_match_frame_percent_f"] = sc["brand_match_frame_percent"].map(to_float)
    sc["overlay_frame_percent_f"] = sc["overlay_frame_percent"].map(to_float)
    sc["cuts_per_minute_f"] = sc["cuts_per_minute"].map(to_float)
    disc_rate = sc["disclosure_ok"].astype(str).eq("1").mean() if "disclosure_ok" in sc else np.nan
    add_kpi_box(s1, 0.3, 1.8, "Avg Brand Tone", f"{np.nanmean(sc['brand_match_frame_percent_f']):.0f}%")
    add_kpi_box(s1, 3.0, 1.8, "Avg Overlay Usage", f"{np.nanmean(sc['overlay_frame_percent_f']):.0f}%")
    add_kpi_box(s1, 5.7, 1.8, "Avg Pacing (cuts/min)", f"{np.nanmean(sc['cuts_per_minute_f']):.1f}")
    add_kpi_box(s1, 8.4, 1.8, "Disclosure Compliance", f"{disc_rate*100:.0f}%")

    # Grade mix by brand
    tmp = sc.copy()
    tmp = ensure_brand_col(tmp)
    mix = tmp.pivot_table(index="brand", columns="overall_grade", values="video_id", aggfunc="count", fill_value=0)
    mix = mix.reindex(sorted(mix.index), fill_value=0)
    fig, ax = plt.subplots()
    mix.plot(kind="bar", stacked=True, ax=ax)
    ax.set_title("Grade Mix by Brand"); ax.set_xlabel(""); ax.set_ylabel("Videos")
    fig_to_picture(s1, fig, 0.4, 3.2, 9.0)

def scatter_and_hist(sc, prs):
    # Scatter: Brand tone vs Overlay
    s2 = prs.slides.add_slide(prs.slide_layouts[5]); add_title(s2, "Brand Tone vs Overlay Usage")
    df2 = sc[["brand_match_frame_percent","overlay_frame_percent"]].copy()
    df2["x"] = df2["brand_match_frame_percent"].map(to_float)
    df2["y"] = df2["overlay_frame_percent"].map(to_float)
    df2 = df2.dropna()
    fig, ax = plt.subplots()
    ax.scatter(df2["x"], df2["y"])
    ax.set_xlabel("Brand-tone adherence (%)"); ax.set_ylabel("Overlay usage (%)"); ax.set_title("Quadrants for action")
    fig_to_picture(s2, fig, 0.5, 1.5, 5.5)

    # Histogram: Pacing
    s3 = prs.slides.add_slide(prs.slide_layouts[5]); add_title(s3, "Pacing (cuts/min) Distribution")
    vals = sc["cuts_per_minute"].map(to_float).dropna()
    fig, ax = plt.subplots()
    ax.hist(vals, bins=12)
    ax.set_xlabel("Cuts per minute"); ax.set_ylabel("Videos")
    fig_to_picture(s3, fig, 0.5, 1.5, 5.5)

def load_thumbs(covers_csv):
    """Return dict video_id -> local image path, preferring saved_path; fall back to downloading cover_url if needed."""
    thumbs = {}
    if not covers_csv: return thumbs
    try:
        df = read_csv_safe(covers_csv)
    except Exception:
        return thumbs
    for _,r in df.iterrows():
        vid = str(r.get("video_id","")).strip()
        saved = str(r.get("saved_path","") or "").strip()
        cover_url = str(r.get("cover_url","") or "").strip()
        if vid and saved and os.path.exists(saved):
            thumbs[vid] = saved
        elif vid and cover_url and HAVE_REQUESTS:
            # try to download to temp next to CSV
            try:
                resp = requests.get(cover_url, timeout=12)
                if resp.status_code==200 and resp.content:
                    out_dir = Path(covers_csv).parent / "covers_dl_cache"
                    out_dir.mkdir(parents=True, exist_ok=True)
                    ext = ".jpg"
                    if "png" in resp.headers.get("Content-Type","").lower(): ext = ".png"
                    outp = out_dir / f"{vid}{ext}"
                    with open(outp,"wb") as f: f.write(resp.content)
                    thumbs[vid] = str(outp)
            except Exception:
                pass
    return thumbs

def add_thumb_card(slide, left, top, width, img_path, title, subtitle, link_url=None, score=None):
    # add picture
    pic = None
    if img_path and os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
        if link_url:
            try:
                pic.click_action.hyperlink.address = link_url
            except Exception:
                pass
    # title box under the image
    box = slide.shapes.add_textbox(Inches(left), Inches(top+width*0.75+0.15), Inches(width), Inches(0.8))
    tf = box.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = (title or "")[:40]
    p.font.size = Pt(12); p.font.bold = True
    if link_url:
        try:
            run = p.runs[0]; run.hyperlink.address = link_url; run.font.color.rgb = RGBColor(0,0,238)
        except Exception:
            pass
    if subtitle or score is not None:
        q = tf.add_paragraph()
        s = subtitle or ""
        if score is not None:
            s = f"{s}  ·  Score {score:.1f}" if s else f"Score {score:.1f}"
        q.text = s; q.font.size = Pt(11)

def top_bottom_grids(sc, prs, thumbs, n=5):
    # rank by composite_score
    def fnum(x):
        try: return float(x)
        except: return -1e9
    sc2 = sc.copy()
    sc2["_rank_score"] = sc2["composite_score"].map(fnum)
    top = sc2.sort_values("_rank_score", ascending=False).head(n)
    bot = sc2.sort_values("_rank_score", ascending=True).head(n)

    # grid positions (3 columns)
    positions = [(0.3,1.4),(3.6,1.4),(6.9,1.4),(0.3,4.2),(3.6,4.2)]
    W = 3.0

    # Top grid
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_title(s, "Top 5 Videos (click to open)")
    for i, (_,r) in enumerate(top.iterrows()):
        if i>=len(positions): break
        vid = str(r.get("video_id",""))
        add_thumb_card(
            s, positions[i][0], positions[i][1], W,
            thumbs.get(vid, None),
            title=f"{r.get('handle','')}/{vid}",
            subtitle=f"{r.get('overall_grade','')} · Brand {r.get('brand_match_frame_percent','')}% · Overlay {r.get('overlay_frame_percent','')}%",
            link_url=r.get("video_url",""),
            score=to_float(r.get("composite_score"))
        )

    # Bottom grid
    s = prs.slides.add_slide(prs.slide_layouts[5]); add_title(s, "Bottom 5 Videos (for coaching)")
    for i, (_,r) in enumerate(bot.iterrows()):
        if i>=len(positions): break
        vid = str(r.get("video_id",""))
        add_thumb_card(
            s, positions[i][0], positions[i][1], W,
            thumbs.get(vid, None),
            title=f"{r.get('handle','')}/{vid}",
            subtitle=f"{r.get('overall_grade','')} · Brand {r.get('brand_match_frame_percent','')}% · Overlay {r.get('overlay_frame_percent','')}%",
            link_url=r.get("video_url",""),
            score=to_float(r.get("composite_score"))
        )

def brand_scorecards(sc, prs, thumbs, n_examples=3):
    brands = list(ensure_brand_col(sc)["brand"].dropna().unique())
    for br in sorted(brands):
        df = sc[sc["brand"]==br].copy()
        if df.empty: continue
        slide = prs.slides.add_slide(prs.slide_layouts[5]); add_title(slide, f"{br} — Scorecard")

        # KPIs
        df["b"] = df["brand_match_frame_percent"].map(to_float)
        df["o"] = df["overlay_frame_percent"].map(to_float)
        df["c"] = df["cuts_per_minute"].map(to_float)
        disc_rate = df["disclosure_ok"].astype(str).eq("1").mean() if "disclosure_ok" in df else np.nan
        add_kpi_box(slide, 0.3, 1.2, "Avg Brand Tone", f"{np.nanmean(df['b']):.0f}%")
        add_kpi_box(slide, 3.0, 1.2, "Avg Overlay", f"{np.nanmean(df['o']):.0f}%")
        add_kpi_box(slide, 5.7, 1.2, "Avg Cuts/min", f"{np.nanmean(df['c']):.1f}")
        add_kpi_box(slide, 8.4, 1.2, "Disclosure OK", f"{disc_rate*100:.0f}%")

        # Top/Bottom examples
        def fnum(x):
            try: return float(x)
            except: return -1e9
        df["_rank"] = df["composite_score"].map(fnum)
        top = df.sort_values("_rank", ascending=False).head(n_examples)
        bot = df.sort_values("_rank", ascending=True).head(n_examples)

        # Layout for examples (thumbnails)
        W = 2.4
        top_positions = [(0.4, 2.8), (3.0, 2.8), (5.6, 2.8)]
        bot_positions = [(0.4, 5.1), (3.0, 5.1), (5.6, 5.1)]

        # Titles
        t1 = slide.shapes.add_textbox(Inches(0.3), Inches(2.5), Inches(4), Inches(0.4)).text_frame
        t1.text = "Top examples"; t1.paragraphs[0].font.bold = True
        t2 = slide.shapes.add_textbox(Inches(0.3), Inches(4.8), Inches(4), Inches(0.4)).text_frame
        t2.text = "Bottom examples (coach)"; t2.paragraphs[0].font.bold = True

        for i, (_,r) in enumerate(top.iterrows()):
            if i>=len(top_positions): break
            vid = str(r.get("video_id",""))
            add_thumb_card(
                slide, top_positions[i][0], top_positions[i][1], W,
                thumbs.get(vid, None),
                title=f"{r.get('handle','')}/{vid}",
                subtitle=f"{r.get('overall_grade','')} · Brand {r.get('brand_match_frame_percent','')}% · Overlay {r.get('overlay_frame_percent','')}%",
                link_url=r.get("video_url",""),
                score=to_float(r.get("composite_score"))
            )

        for i, (_,r) in enumerate(bot.iterrows()):
            if i>=len(bot_positions): break
            vid = str(r.get("video_id",""))
            add_thumb_card(
                slide, bot_positions[i][0], bot_positions[i][1], W,
                thumbs.get(vid, None),
                title=f"{r.get('handle','')}/{vid}",
                subtitle=f"{r.get('overall_grade','')} · Brand {r.get('brand_match_frame_percent','')}% · Overlay {r.get('overlay_frame_percent','')}%",
                link_url=r.get("video_url",""),
                score=to_float(r.get("composite_score"))
            )

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--scoring", required=True, help="Path to scoring_summary.csv from v3")
    ap.add_argument("--covers", default=None, help="Optional covers_manifest.csv for thumbnails")
    ap.add_argument("--out", default="Exec_Summary.pptx")
    args = ap.parse_args()

    sc = read_csv_safe(args.scoring)
    # Ensure required numeric columns are strings -> floats only where necessary
    for c in ["brand_match_frame_percent","overlay_frame_percent","cuts_per_minute","disclosure_ok","composite_score"]:
        if c in sc.columns: sc[c] = sc[c].astype(str)

    prs = Presentation()

    brand_mix_chart(sc, prs)
    scatter_and_hist(sc, prs)

    thumbs = load_thumbs(args.covers) if args.covers else {}
    top_bottom_grids(sc, prs, thumbs, n=5)
    brand_scorecards(sc, prs, thumbs, n_examples=3)

    prs.save(args.out)
    print(f"Saved: {args.out}")

if __name__ == "__main__":
    main()
